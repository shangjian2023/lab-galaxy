"""Document parser — extract text from Word / PDF / PPT files."""

import io
import subprocess
import tempfile
from pathlib import Path

from app.registries.parsers import parser_registry


@parser_registry.register(".docx")
def parse_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


@parser_registry.register(".doc")
def parse_doc(data: bytes) -> str:
    """Parse legacy .doc files via libreoffice headless conversion."""
    try:
        with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
            tmp.write(data)
            tmp_path = tmp.name
        out_dir = Path(tmp_path).parent
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt", "--outdir", str(out_dir), tmp_path],
            capture_output=True, timeout=30,
        )
        txt_path = Path(tmp_path).with_suffix(".txt")
        if txt_path.exists():
            text = txt_path.read_text(errors="replace")
            Path(tmp_path).unlink(missing_ok=True)
            txt_path.unlink(missing_ok=True)
            return "\n".join(line for line in text.splitlines() if line.strip())
        raise RuntimeError("libreoffice conversion produced no output")
    except (FileNotFoundError, subprocess.TimeoutExpired, RuntimeError):
        # Fallback: try reading raw bytes as text (partial support)
        return data.decode("utf-8", errors="replace")


@parser_registry.register(".pdf")
def parse_pdf(data: bytes) -> str:
    import fitz

    text_parts: list[str] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for page in doc:
            text_parts.append(page.get_text())
    return "\n".join(text_parts)


@parser_registry.register(".pptx")
@parser_registry.register(".ppt")
def parse_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    text_parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        text_parts.append(para.text)
    return "\n".join(text_parts)


def parse_file(data: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    try:
        parser_fn = parser_registry.get(ext)
    except KeyError:
        raise ValueError(f"不支持的文件格式: {ext}")
    return parser_fn(data)
